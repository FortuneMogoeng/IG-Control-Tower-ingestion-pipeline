"""
parse_pptx.py — PowerPoint (.pptx) parser for AI Control Tower ingestion.
Extracts title + body text per slide and emits provenance-tagged fragments
as JSON Lines (one per slide, skipping blank slides).

Usage:
    python parse_pptx.py <path/to/file.pptx>

Output (stdout, one JSON object per line):
    {
        "fragment_id": "uuid4",
        "source_file": "filename.pptx",
        "file_type": "pptx",
        "location": "Slide 4: 'Current Process Overview'",
        "text": "Current Process Overview\nManual data entry takes 3 days per week...",
        "char_count": 198
    }

Requires: python-pptx  (pip install python-pptx)
"""

import json
import sys
import uuid
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")


MIN_CHARS = 30


def slide_text(slide) -> tuple[str, str]:
    """Return (title, body_text) for a slide. Title may be empty string."""
    title = ""
    body_parts = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        # Detect title placeholder
        if shape.shape_id == 1 or (
            hasattr(shape, "placeholder_format")
            and shape.placeholder_format is not None
            and shape.placeholder_format.idx == 0
        ):
            title = text
        else:
            body_parts.append(text)

    return title, "\n".join(body_parts)


def parse(path: str) -> list[dict]:
    source = Path(path)
    prs = Presentation(source)
    fragments: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title, body = slide_text(slide)
        combined = f"{title}\n{body}".strip() if title else body.strip()

        if len(combined) < MIN_CHARS:
            continue

        location = f"Slide {slide_num}"
        if title:
            location += f": '{title}'"

        fragments.append({
            "fragment_id": str(uuid.uuid4()),
            "source_file": source.name,
            "file_type": "pptx",
            "location": location,
            "text": combined,
            "char_count": len(combined),
        })

    return fragments


def main():
    if len(sys.argv) < 2:
        sys.exit("Usage: python parse_pptx.py <file.pptx>")

    for fragment in parse(sys.argv[1]):
        print(json.dumps(fragment))


if __name__ == "__main__":
    main()
